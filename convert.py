import pdfplumber
import pandas as pd
from docx import Document
from pdf2docx import Converter
from pptx import Presentation
from pptx.util import Inches
import fitz  # PyMuPDF
from pypdf import PdfReader, PdfWriter   # switched to modern package
import pytesseract
from PIL import Image
import os
import sys

# ✅ Set Tesseract path (update if installed elsewhere)
pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"

# ------------------------------
# Excel Converter
# ------------------------------
def pdf_to_excel(pdf_path, excel_path="output.xlsx"):
    all_tables = []
    with pdfplumber.open(pdf_path) as pdf:
        total_pages = len(pdf.pages)
        for page_num, page in enumerate(pdf.pages, start=1):
            print(f"Processing page {page_num}/{total_pages} for tables...")
            tables = page.extract_tables()
            if not tables:
                continue
            for t_index, table in enumerate(tables, start=1):
                if table:
                    headers = table[0]
                    headers = pd.io.parsers.ParserBase({'names': headers})._maybe_dedup_names(headers)
                    df = pd.DataFrame(table[1:], columns=headers)
                    all_tables.append(df)
                    print(f"  ✔ Found table {t_index} on page {page_num}")
    if all_tables:
        final_df = pd.concat(all_tables, ignore_index=True)
        final_df.to_excel(excel_path, index=False)
        print(f"✅ Saved tables to {excel_path}")
    else:
        # Fallback: dump text into Excel
        with pdfplumber.open(pdf_path) as pdf:
            text = []
            total_pages = len(pdf.pages)
            for page_num, page in enumerate(pdf.pages, start=1):
                print(f"Processing page {page_num}/{total_pages} for text fallback...")
                page_text = page.extract_text()
                text.append(page_text if page_text else "")
        df = pd.DataFrame({"Text": text})
        df.to_excel(excel_path, index=False)
        print(f"⚠ No tables found, saved raw text to {excel_path}")

# ------------------------------
# Word Converter (layout-preserved)
# ------------------------------
def pdf_to_word(pdf_path, word_path="output.docx"):
    print(f"Converting {pdf_path} to Word (layout preserved)...")
    cv = Converter(pdf_path)
    cv.convert(word_path, start=0, end=None)  # all pages
    cv.close()
    print(f"✅ Saved to {word_path}")

# ------------------------------
# PPT Converter (pages as images)
# ------------------------------
def pdf_to_ppt(pdf_path, ppt_path="output.pptx"):
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]  # completely blank slide
    doc = fitz.open(pdf_path)
    total_pages = len(doc)

    for page_num in range(total_pages):
        print(f"Rendering page {page_num+1}/{total_pages} as slide...")
        slide = prs.slides.add_slide(blank_slide_layout)
        pix = doc[page_num].get_pixmap()
        img_path = f"temp_page_{page_num}.png"
        pix.save(img_path)
        slide.shapes.add_picture(img_path, Inches(0), Inches(0), Inches(10), Inches(7.5))
        os.remove(img_path)

    prs.save(ppt_path)
    print(f"✅ Saved to {ppt_path}")

# ------------------------------
# Compress PDF (Fixed)
# ------------------------------
def compress_pdf(pdf_path, output_path="compressed.pdf", dpi=100, quality=60):
    doc = fitz.open(pdf_path)
    new_doc = fitz.open()
    for page_num in range(len(doc)):
        print(f"Compressing page {page_num+1}/{len(doc)}...")
        pix = doc[page_num].get_pixmap(matrix=fitz.Matrix(dpi/72, dpi/72))
        img_bytes = pix.tobytes("jpg")  # ✅ use jpeg instead of pdf
        img_pdf = fitz.open("pdf", img_bytes)
        new_doc.insert_pdf(img_pdf)
    new_doc.save(output_path, deflate=True)
    print(f"✅ Compressed PDF saved → {output_path}")

# ------------------------------
# Password Protect / Unprotect
# ------------------------------
def protect_pdf(pdf_path, output_path="protected.pdf", password="1234"):
    reader = PdfReader(pdf_path)
    writer = PdfWriter()
    for page in reader.pages:
        writer.add_page(page)
    writer.encrypt(password)
    with open(output_path, "wb") as f:
        writer.write(f)
    print(f"🔒 Protected PDF saved → {output_path}")

def unprotect_pdf(pdf_path, output_path="unprotected.pdf", password="1234"):
    reader = PdfReader(pdf_path)
    if reader.is_encrypted:
        try:
            reader.decrypt(password)
        except:
            print("❌ Wrong password. Cannot decrypt PDF.")
            return
    writer = PdfWriter()
    for page in reader.pages:
        writer.add_page(page)
    with open(output_path, "wb") as f:
        writer.write(f)
    print(f"🔓 Unprotected PDF saved → {output_path}")

# ------------------------------
# OCR for Scanned PDFs (Fixed)
# ------------------------------
def ocr_pdf(pdf_path, output_path="ocr_output.docx"):
    doc = Document()
    with fitz.open(pdf_path) as pdf:
        total_pages = len(pdf)
        for page_num in range(total_pages):
            print(f"OCR on page {page_num+1}/{total_pages}...")
            pix = pdf[page_num].get_pixmap()
            img_path = f"temp_page_{page_num}.png"
            pix.save(img_path)
            text = pytesseract.image_to_string(Image.open(img_path), lang="eng")
            os.remove(img_path)
            doc.add_heading(f"Page {page_num+1}", level=2)
            doc.add_paragraph(text)
    doc.save(output_path)
    print(f"✅ OCR Word file saved → {output_path}")

# ------------------------------
# Auto-Detection
# ------------------------------
def detect_mode(pdf_path):
    text_count, table_count = 0, 0
    with pdfplumber.open(pdf_path) as pdf:
        total_pages = len(pdf.pages)
        for page_num, page in enumerate(pdf.pages, start=1):
            print(f"Analyzing page {page_num}/{total_pages}...")
            text = page.extract_text()
            if text:
                text_count += len(text.split())
            tables = page.extract_tables()
            if tables:
                table_count += len(tables)

    if table_count > text_count / 50:
        return "excel"
    elif text_count > 200:
        return "word"
    else:
        return "ppt"

# ------------------------------
# Main CLI
# ------------------------------
if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python convert.py <input.pdf> [--output=auto|excel|word|ppt|compress|protect|unprotect|ocr] [--password=xxxx]")
        sys.exit(1)

    pdf_file = sys.argv[1]
    if not os.path.exists(pdf_file):
        print(f"❌ File not found: {pdf_file}")
        sys.exit(1)

    # Default mode
    output_mode = "auto"
    password = None

    # Parse args
    for arg in sys.argv[2:]:
        if arg.startswith("--output"):
            parts = arg.split("=")
            if len(parts) == 2:
                output_mode = parts[1].strip().lower()
        elif arg.startswith("--password"):
            parts = arg.split("=")
            if len(parts) == 2:
                password = parts[1].strip()

    if output_mode == "auto":
        output_mode = detect_mode(pdf_file)
        print(f"🔍 Auto-detected best format: {output_mode}")

    if output_mode == "excel":
        pdf_to_excel(pdf_file, "output.xlsx")
    elif output_mode == "word":
        pdf_to_word(pdf_file, "output.docx")
    elif output_mode == "ppt":
        pdf_to_ppt(pdf_file, "output.pptx")
    elif output_mode == "compress":
        compress_pdf(pdf_file, "compressed.pdf")
    elif output_mode == "protect":
        if not password:
            print("❌ Please provide a password with --password=xxxx")
        else:
            protect_pdf(pdf_file, "protected.pdf", password)
    elif output_mode == "unprotect":
        if not password:
            print("❌ Please provide the current password with --password=xxxx")
        else:
            unprotect_pdf(pdf_file, "unprotected.pdf", password)
    elif output_mode == "ocr":
        ocr_pdf(pdf_file, "ocr_output.docx")
    else:
        print("❌ Unknown output mode. Use --output=auto|excel|word|ppt|compress|protect|unprotect|ocr")
